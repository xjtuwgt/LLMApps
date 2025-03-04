import json
from pptx import Presentation

# Sample JSON Data
json_data = '''
{
    "title": "Project Overview",
    "slides": [
        {
            "heading": "Introduction",
            "content": "This presentation gives an overview of our project."
        },
        {
            "heading": "Objectives",
            "content": "1. Improve efficiency\\n2. Enhance user experience\\n3. Increase scalability"
        },
        {
            "heading": "Conclusion",
            "content": "Thank you for your attention!"
        }
    ]
}
'''

def json_to_pptx(json_data):
    data = json.loads(json_data)

    # Create a PowerPoint Presentation
    presentation = Presentation()

    # Add Title Slide
    title_slide_layout = presentation.slide_layouts[0]  # Title slide layout
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = data["title"]
    subtitle.text = "Generated from JSON"

    # Add Content Slides
    for slide_data in data["slides"]:
        heading = slide_data["heading"]
        content_items = slide_data["content"]

        # Ensure content is in list format for uniform processing
        if isinstance(content_items, str):  
            content_items = [content_items]

        # Create multiple slides if multiple content items exist
        for content_item in content_items:
            slide_layout = presentation.slide_layouts[1]  # Title and content layout
            slide = presentation.slides.add_slide(slide_layout)

            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = heading  # Same heading for all slides in this category
            content.text = content_item  # Different content for each slide

    return presentation


# Save the Presentation
pptx_filename = "output_presentation.pptx"
presentation = json_to_pptx(json_data)
presentation.save(pptx_filename)

print(f"Presentation saved as {pptx_filename}")